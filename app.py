import io
import json
import os
import re
import time
import base64
import logging
from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime
import asyncio
from contextlib import asynccontextmanager

from fastapi import FastAPI, UploadFile, Form, HTTPException, Request
from fastapi.responses import HTMLResponse, StreamingResponse, FileResponse, Response
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from fastapi.middleware.gzip import GZipMiddleware

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Optional AI SDK imports with graceful fallback
try:
    from openai import OpenAI
    HAS_OPENAI = True
except ImportError:
    OpenAI = None
    HAS_OPENAI = False

try:
    import anthropic
    HAS_ANTHROPIC = True
except ImportError:
    anthropic = None
    HAS_ANTHROPIC = False

try:
    import google.generativeai as genai
    HAS_GEMINI = True
except ImportError:
    genai = None
    HAS_GEMINI = False

try:
    import requests
    HAS_REQUESTS = True
except ImportError:
    requests = None
    HAS_REQUESTS = False

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Configuration
class Config:
    # Application settings
    APP_NAME = "SlideForge AI"
    APP_VERSION = "1.0.0"
    DEBUG = os.getenv("DEBUG", "false").lower() == "true"
    
    # Limits
    MAX_TEXT_CHARS = 100_000
    MIN_SLIDES = 1
    MAX_SLIDES = 50
    MAX_TEMPLATE_BYTES = 50 * 1024 * 1024  # 50 MB
    MAX_CONTENT_LENGTH = 100 * 1024 * 1024  # 100 MB
    
    # Default models
    DEFAULT_MODELS = {
        "openai": "gpt-4o-mini",
        "anthropic": "claude-3-5-sonnet-latest",
        "gemini": "gemini-2.5-flash",
        "aipipe": "gpt-4o-mini"
    }
    
    # AI Pipe endpoint
    AIPIPE_ENDPOINT = "https://aipipe.org/openrouter/v1/chat/completions"
    
    # Request timeouts
    AI_REQUEST_TIMEOUT = 120
    MAX_RETRIES = 3
    RETRY_DELAY = 1.0

config = Config()

# Application lifespan
@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info(f"Starting {config.APP_NAME} v{config.APP_VERSION}")
    logger.info(f"Available AI providers: OpenAI={HAS_OPENAI}, Anthropic={HAS_ANTHROPIC}, Gemini={HAS_GEMINI}")
    yield
    logger.info("Shutting down SlideForge AI")

# Initialize FastAPI app
app = FastAPI(
    title=config.APP_NAME,
    version=config.APP_VERSION,
    description="AI-powered presentation generator with multi-provider support",
    lifespan=lifespan
)

from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
import os

# Serve frontend
@app.get("/", include_in_schema=False)
async def serve_index():
    return FileResponse("Index.html")

# Serve favicon (optional, since you already uploaded favicon.ico)
@app.get("/favicon.ico", include_in_schema=False)
async def favicon():
    return FileResponse("favicon.ico")

# Add middleware
app.add_middleware(GZipMiddleware, minimum_size=1000)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Serve static files and frontend
@app.get("/", response_class=HTMLResponse)
async def serve_frontend():
    """Serve the SlideForge AI frontend"""
    html_path = os.path.join(os.path.dirname(__file__), "index.html")
    try:
        with open(html_path, "r", encoding="utf-8") as f:
            content = f.read()
            logger.info("Served frontend successfully")
            return HTMLResponse(content=content)
    except FileNotFoundError:
        logger.error("Frontend HTML file not found")
        return HTMLResponse(
            content="""
            <html>
                <head><title>SlideForge AI - File Not Found</title></head>
                <body style="font-family: system-ui; padding: 2rem; text-align: center;">
                    <h1>SlideForge AI Backend</h1>
                    <p>Frontend not found. Please ensure index.html is in the same directory as this script.</p>
                    <p>Backend is running successfully on: <code>/generate</code></p>
                </body>
            </html>
            """,
            status_code=404,
        )

# Favicon with fallback
@app.get("/favicon.ico", include_in_schema=False)
async def favicon():
    """Serve favicon or return empty response"""
    favicon_path = "favicon.ico"
    if os.path.exists(favicon_path):
        return FileResponse(favicon_path, media_type="image/x-icon")
    
    # Return a tiny transparent PNG as fallback
    tiny_png = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO3n+9QAAAAASUVORK5CYII="
    )
    return Response(content=tiny_png, media_type="image/png")

# Health check endpoint
@app.get("/health")
async def health_check():
    """Health check endpoint for monitoring"""
    return {
        "status": "healthy",
        "app": config.APP_NAME,
        "version": config.APP_VERSION,
        "timestamp": datetime.utcnow().isoformat(),
        "providers": {
            "openai": HAS_OPENAI,
            "anthropic": HAS_ANTHROPIC,
            "gemini": HAS_GEMINI,
            "requests": HAS_REQUESTS
        }
    }

# AI Provider Classes
class AIProvider:
    """Base class for AI providers"""
    
    def __init__(self, api_key: str, model: str):
        self.api_key = api_key.strip()
        self.model = model.strip()
        self.timeout = config.AI_REQUEST_TIMEOUT
    
    async def generate_slides(self, prompt: str) -> Dict[str, Any]:
        raise NotImplementedError
    
    def _create_slide_prompt(self, text: str, guidance: str, target_slides: Optional[int]) -> str:
        """Create a standardized prompt for slide generation"""
        slide_count_instruction = (
             f"You MUST create exactly {target_slides} slides. "
             f"Do not create placeholder or generic 'Summary' slides."
             f"If you cannot generate enough meaningful content, expand on subtopics instead"
             f"Do not return fewer or more. "
             f"If necessary, split or expand content so that there are exactly {target_slides} slides."
             if target_slides else 
             f"Create between {config.MIN_SLIDES} and {config.MAX_SLIDES} slides."
        )

        
        return f"""You are an expert presentation designer. Create a slide deck from the provided content.

REQUIREMENTS:
- {slide_count_instruction}
- Each slide must have a title (max 80 characters)
- Each slide should have 3-6 bullet points (max 120 characters each)
- Focus on clear, actionable content
- Style guidance: {guidance or 'professional and engaging'}

OUTPUT FORMAT:
Return ONLY a JSON object with this exact structure:
{{
    "slides": [
        {{
            "title": "Slide title here",
            "bullets": ["Bullet point 1", "Bullet point 2", "Bullet point 3"]
        }}
    ]
}}

CONTENT TO CONVERT:
{text}

Return only valid JSON, no explanations or markdown formatting."""

class OpenAIProvider(AIProvider):
    """OpenAI API provider"""
    
    async def generate_slides(self, prompt: str) -> Dict[str, Any]:
        if not HAS_OPENAI:
            raise HTTPException(500, "OpenAI package not installed")
        
        client = OpenAI(api_key=self.api_key, timeout=self.timeout)
        
        try:
            response = client.chat.completions.create(
                model=self.model,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=2048,
                response_format={"type": "json_object"}
            )
            
            content = response.choices[0].message.content
            return self._parse_json_response(content)
            
        except Exception as e:
            logger.error(f"OpenAI API error: {e}")
            raise HTTPException(500, f"OpenAI API error: {str(e)}")

class AnthropicProvider(AIProvider):
    """Anthropic Claude API provider"""
    
    async def generate_slides(self, prompt: str) -> Dict[str, Any]:
        if not HAS_ANTHROPIC:
            raise HTTPException(500, "Anthropic package not installed")
        
        client = anthropic.Anthropic(api_key=self.api_key, timeout=self.timeout)
        
        try:
            message = client.messages.create(
                model=self.model,
                max_tokens=2048,
                temperature=0.3,
                system="You are a presentation expert. Return only valid JSON responses.",
                messages=[{"role": "user", "content": prompt}]
            )
            
            content = "".join([
                block.text for block in message.content 
                if hasattr(block, 'text')
            ])
            
            return self._parse_json_response(content)
            
        except Exception as e:
            logger.error(f"Anthropic API error: {e}")
            raise HTTPException(500, f"Anthropic API error: {str(e)}")

class GeminiProvider(AIProvider):
    """Google Gemini API provider"""
    
    async def generate_slides(self, prompt: str) -> Dict[str, Any]:
        if not HAS_GEMINI:
            raise HTTPException(500, "Google Generative AI package not installed")
        
        try:
            genai.configure(api_key=self.api_key)
            model = genai.GenerativeModel(self.model)
            
            # Add safety settings to prevent blocking
            safety_settings = [
                {
                    "category": "HARM_CATEGORY_HARASSMENT",
                    "threshold": "BLOCK_NONE"
                },
                {
                    "category": "HARM_CATEGORY_HATE_SPEECH",
                    "threshold": "BLOCK_NONE"
                },
                {
                    "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
                    "threshold": "BLOCK_NONE"
                },
                {
                    "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
                    "threshold": "BLOCK_NONE"
                }
            ]
            
            response = model.generate_content(
                prompt,
                generation_config=genai.types.GenerationConfig(
                    temperature=0.3,
                    max_output_tokens=2048,
                    response_mime_type="application/json"  # Request JSON response
                ),
                safety_settings=safety_settings
            )
            
            # Check if response has valid parts
            if not response.parts:
                # Check finish reason
                if response.candidates and response.candidates[0].finish_reason:
                    finish_reason = response.candidates[0].finish_reason
                    if finish_reason == 2:  # SAFETY
                        logger.warning("Gemini blocked response due to safety filters")
                        # Return a fallback response
                        return {
                            "slides": [{
                                "title": "Content Generation Notice",
                                "bullets": [
                                    "Content was filtered by safety systems",
                                    "Please try rephrasing your input",
                                    "Ensure content follows guidelines"
                                ]
                            }]
                        }
                    elif finish_reason == 3:  # RECITATION
                        logger.warning("Gemini blocked response due to recitation")
                        return {
                            "slides": [{
                                "title": "Content Generation Notice",
                                "bullets": ["Content too similar to training data", "Please provide more unique input"]
                            }]
                        }
                
                # Generic fallback
                raise HTTPException(500, "Gemini returned empty response")
            
            # Extract text from response
            text_content = response.text if hasattr(response, 'text') else response.parts[0].text
            return self._parse_json_response(text_content)
            
        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Gemini API error: {e}")
            raise HTTPException(500, f"Gemini API error: {str(e)}")

class AIPipeProvider(AIProvider):
    """AI Pipe (OpenRouter) API provider"""
    
    async def generate_slides(self, prompt: str) -> Dict[str, Any]:
        if not HAS_REQUESTS:
            raise HTTPException(500, "Requests package not installed")
        
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.api_key}",
        }
        
        payload = {
            "model": self.model,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.3,
            "max_tokens": 2048,
        }
        
        try:
            response = requests.post(
                config.AIPIPE_ENDPOINT,
                json=payload,
                headers=headers,
                timeout=self.timeout
            )
            
            if not response.ok:
                raise HTTPException(
                    response.status_code, 
                    f"AI Pipe API error: {response.text}"
                )
            
            data = response.json()
            content = data["choices"][0]["message"]["content"]
            
            return self._parse_json_response(content)
            
        except requests.RequestException as e:
            logger.error(f"AI Pipe API error: {e}")
            raise HTTPException(500, f"AI Pipe API error: {str(e)}")
        except KeyError as e:
            logger.error(f"AI Pipe response format error: {e}")
            raise HTTPException(500, "Invalid AI Pipe API response format")

# Add JSON parsing method to base class
def _parse_json_response(self, content: str) -> Dict[str, Any]:
    """Parse JSON response with error handling"""
    if not content or not content.strip():
        return {"slides": [{"title": "Error", "bullets": ["No content generated"]}]}
    
    try:
        # Try direct JSON parsing first
        data = json.loads(content.strip())
        if isinstance(data, dict) and "slides" in data:
            return data
    except json.JSONDecodeError:
        pass
    
    # Try to extract JSON from markdown or other formatting
    json_match = re.search(r'\{.*\}', content.strip(), re.DOTALL)
    if json_match:
        try:
            data = json.loads(json_match.group())
            if isinstance(data, dict) and "slides" in data:
                return data
        except json.JSONDecodeError:
            pass
    
    # Fallback: create a single slide with the content
    logger.warning(f"Failed to parse JSON response, creating fallback slide")
    return {
        "slides": [{
            "title": "Generated Content", 
            "bullets": [content[:120] + "..." if len(content) > 120 else content]
        }]
    }

# Add the method to base class
AIProvider._parse_json_response = _parse_json_response

# Provider factory
def create_provider(provider_name: str, api_key: str, model: str) -> AIProvider:
    """Create AI provider instance"""
    provider_map = {
        "openai": OpenAIProvider,
        "anthropic": AnthropicProvider,
        "gemini": GeminiProvider,
        "aipipe": AIPipeProvider,
    }
    
    if provider_name not in provider_map:
        raise HTTPException(400, f"Unsupported provider: {provider_name}")
    
    return provider_map[provider_name](api_key, model)

# Slide processing utilities
class SlideProcessor:
    """Utilities for processing and validating slide data"""
    
    @staticmethod
    def normalize_slides(slides_data: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Normalize slide data format"""
        normalized = []
        
        for slide in slides_data:
            if not isinstance(slide, dict):
                continue
                
            title = str(slide.get("title", "")).strip()[:80] or "Untitled Slide"
            bullets = slide.get("bullets", [])
            
            if not isinstance(bullets, list):
                bullets = []
            
            # Clean and limit bullets
            clean_bullets = []
            for bullet in bullets[:6]:  # Max 6 bullets
                if bullet and str(bullet).strip():
                    clean_bullets.append(str(bullet).strip()[:120])
            
            normalized.append({
                "title": title,
                "bullets": clean_bullets
            })
        
        return normalized
    
    @staticmethod
    def adjust_slide_count(slides: List[Dict[str, Any]], target_count: Optional[int]) -> List[Dict[str, Any]]:
        """Adjust slide count to meet target"""
        if not target_count:
            return slides
        
        target_count = max(config.MIN_SLIDES, min(config.MAX_SLIDES, target_count))
        
        if len(slides) == target_count:
            return slides
        
        if len(slides) > target_count:
            # Truncate excess slides
            return slides[:target_count]
        
        # Need more slides - split content or add summary slides
        while len(slides) < target_count:
            # Find the slide with the most bullets to split
            max_bullets = 0
            split_idx = 0
            
            for i, slide in enumerate(slides):
                if len(slide["bullets"]) > max_bullets:
                    max_bullets = len(slide["bullets"])
                    split_idx = i
            
            if max_bullets > 2:
                # Split the slide
                original = slides[split_idx]
                mid_point = len(original["bullets"]) // 2
                
                slides[split_idx] = {
                    "title": original["title"],
                    "bullets": original["bullets"][:mid_point]
                }
                
                slides.insert(split_idx + 1, {
                    "title": original["title"] + " (continued)",
                    "bullets": original["bullets"][mid_point:]
                })
        
        return slides[:target_count]

# PowerPoint generation utilities
class PPTXBuilder:
    """PowerPoint presentation builder"""
    
    def __init__(self, template_bytes: Optional[bytes] = None):
        self.template_bytes = template_bytes
        self.prs = None
        self._initialize_presentation()
    
    def _initialize_presentation(self):
        """Initialize presentation from template or blank"""
        if self.template_bytes:
            try:
                self.prs = Presentation(io.BytesIO(self.template_bytes))
                logger.info("Loaded presentation from template")
            except Exception as e:
                logger.warning(f"Failed to load template: {e}, using blank presentation")
                self.prs = Presentation()
        else:
            self.prs = Presentation()
        
        # Clear existing slides if using template
        if self.template_bytes and len(self.prs.slides) > 0:
            self._clear_slides()
    
    def _clear_slides(self):
        """Safely remove all slides from presentation"""
        xml_slides = self.prs.slides._sldIdLst
        for sld_id in list(xml_slides):
            r_id = sld_id.rId
            self.prs.part.drop_rel(r_id)
            xml_slides.remove(sld_id)
    
    def _get_layout(self) -> int:
        """Get the best layout index for content slides"""
        # Look for title and content layout
        for i, layout in enumerate(self.prs.slide_layouts):
            try:
                placeholders = list(layout.placeholders)
                has_title = any(ph.placeholder_format.type == 1 for ph in placeholders)  # Title
                has_content = any(ph.placeholder_format.type in [2, 7] for ph in placeholders)  # Body/Content
                
                if has_title and has_content:
                    return i
            except:
                continue
        
        # Fallback to layout 1 or 0
        return min(1, len(self.prs.slide_layouts) - 1)
    
    def add_slide(self, title: str, bullets: List[str]) -> bool:
        """Add a slide with title and bullets"""
        try:
            layout_idx = self._get_layout()
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
            
            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = title
                # Style the title
                title_shape = slide.shapes.title
                title_shape.text_frame.paragraphs[0].font.size = Pt(28)
                title_shape.text_frame.paragraphs[0].font.bold = True
            
            # Add bullets to content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                try:
                    if shape.placeholder_format.type in [2, 7]:  # Body or Content
                        content_placeholder = shape
                        break
                except:
                    continue
            
            if content_placeholder and bullets:
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                # Add first bullet
                if bullets:
                    p = text_frame.paragraphs[0]
                    p.text = bullets[0]
                    p.level = 0
                    p.font.size = Pt(18)
                    
                    # Add remaining bullets
                    for bullet in bullets[1:]:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(18)
            
            elif bullets:
                # No content placeholder, create text box
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(5)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                
                if bullets:
                    p = text_frame.paragraphs[0]
                    p.text = bullets[0]
                    p.font.size = Pt(18)
                    
                    for bullet in bullets[1:]:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(18)
            
            return True
            
        except Exception as e:
            logger.error(f"Error adding slide: {e}")
            return False
    
    def save(self) -> bytes:
        """Save presentation to bytes"""
        output = io.BytesIO()
        self.prs.save(output)
        return output.getvalue()

# Main generation endpoint
@app.post("/generate")
async def generate_presentation(
    text: str = Form(...),
    guidance: Optional[str] = Form(None),
    provider: str = Form(...),
    api_key: str = Form(...),
    model: Optional[str] = Form(None),
    num_slides: Optional[int] = Form(None),
    reuse_images: bool = Form(False),
    template: Optional[UploadFile] = None,
    request: Request = None
):
    """Generate PowerPoint presentation using AI"""
    start_time = time.time()
    client_ip = request.client.host if request else "unknown"
    
    logger.info(f"Generation request from {client_ip}: provider={provider}, model={model}")
    
    # Validate inputs
    if not text or len(text.strip()) < 10:
        raise HTTPException(400, "Text content must be at least 10 characters")
    
    if not api_key or len(api_key.strip()) < 6:
        raise HTTPException(400, "Valid API key is required")
    
    text = text.strip()[:config.MAX_TEXT_CHARS]
    guidance = (guidance or "").strip()
    provider = provider.lower().strip()
    
    # Set default model if not provided
    if not model:
        model = config.DEFAULT_MODELS.get(provider, "gpt-4o-mini")
    
    # Validate slide count
    target_slides = None
    if num_slides is not None:
        target_slides = max(config.MIN_SLIDES, min(config.MAX_SLIDES, int(num_slides)))
    
    # Handle template upload
    template_bytes = None
    if template and template.filename:
        if not template.filename.lower().endswith(('.pptx', '.potx')):
            raise HTTPException(400, "Template must be a .pptx or .potx file")
        
        template_bytes = await template.read()
        if len(template_bytes) > config.MAX_TEMPLATE_BYTES:
            raise HTTPException(400, f"Template too large (max {config.MAX_TEMPLATE_BYTES // (1024*1024)}MB)")
        
        if len(template_bytes) < 1024:
            raise HTTPException(400, "Template file appears to be empty or corrupted")
    
    try:
        # Create AI provider and generate slides
        ai_provider = create_provider(provider, api_key, model)
        prompt = ai_provider._create_slide_prompt(text, guidance, target_slides)
        
        logger.info(f"Generating slides with {provider}/{model}")
        slides_data = await ai_provider.generate_slides(prompt)
        
        # Process slides
        raw_slides = slides_data.get("slides", [])
        if not raw_slides:
            raise HTTPException(500, "AI provider returned no slides")
        
        # Normalize and adjust slide count
        normalized_slides = SlideProcessor.normalize_slides(raw_slides)
        final_slides = SlideProcessor.adjust_slide_count(normalized_slides, target_slides)
        
        logger.info(f"Generated {len(final_slides)} slides")
        
        # Build PowerPoint
        builder = PPTXBuilder(template_bytes)
        
        for slide_data in final_slides:
            success = builder.add_slide(slide_data["title"], slide_data["bullets"])
            if not success:
                logger.warning(f"Failed to add slide: {slide_data['title']}")
        
        # Generate final file
        pptx_bytes = builder.save()
        
        # Log success
        generation_time = time.time() - start_time
        logger.info(f"Successfully generated presentation in {generation_time:.2f}s "
                   f"(size: {len(pptx_bytes)} bytes)")
        
        # Return file
        headers = {
            "Content-Disposition": 'attachment; filename="slideforge_presentation.pptx"',
            "Content-Type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "X-Generation-Time": str(generation_time),
            "X-Slides-Count": str(len(final_slides))
        }
        
        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers=headers
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Generation failed: {e}", exc_info=True)
        raise HTTPException(500, f"Generation failed: {str(e)}")

# Error handlers
@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    """Custom HTTP exception handler with detailed logging"""
    client_ip = request.client.host
    logger.warning(f"HTTP {exc.status_code} from {client_ip}: {exc.detail}")
    
    return Response(
        content=json.dumps({
            "error": exc.detail,
            "status_code": exc.status_code,
            "timestamp": datetime.utcnow().isoformat()
        }),
        status_code=exc.status_code,
        media_type="application/json"
    )

@app.exception_handler(Exception)
async def general_exception_handler(request: Request, exc: Exception):
    """Handle unexpected errors"""
    client_ip = request.client.host
    logger.error(f"Unexpected error from {client_ip}: {exc}", exc_info=True)
    
    return Response(
        content=json.dumps({
            "error": "Internal server error",
            "status_code": 500,
            "timestamp": datetime.utcnow().isoformat()
        }),
        status_code=500,
        media_type="application/json"
    )

# Development server
if __name__ == "__main__":
    import uvicorn
    
    logger.info("Starting SlideForge AI development server...")
    uvicorn.run(
        "app:app",
        host="0.0.0.0",
        port=8000,
        reload=config.DEBUG,
        log_level="info" if not config.DEBUG else "debug"

    )
